#!/usr/bin/env python3
"""
HTML to PPTX Converter
Converts HTML files (with CSS, JS, and design) into PowerPoint presentations
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.chrome.service import Service
from webdriver_manager.chrome import ChromeDriverManager
from PIL import Image
import io
import time

class HTMLtoPPTXConverter:
    def __init__(self, width=1920, height=1080):
        self.width = width
        self.height = height
        self.driver = None
        
    def setup_browser(self):
        """Initialize headless Chrome browser"""
        chrome_options = Options()
        chrome_options.add_argument('--headless')
        chrome_options.add_argument('--no-sandbox')
        chrome_options.add_argument('--disable-dev-shm-usage')
        chrome_options.add_argument(f'--window-size={self.width},{self.height}')
        
        service = Service(ChromeDriverManager().install())
        self.driver = webdriver.Chrome(service=service, options=chrome_options)
        
    def capture_html(self, html_path, wait_time=2):
        """Capture rendered HTML as screenshot"""
        if not self.driver:
            self.setup_browser()
            
        # Load HTML file
        file_url = f'file://{os.path.abspath(html_path)}'
        self.driver.get(file_url)
        
        # Wait for JS/CSS to load
        time.sleep(wait_time)
        
        # Get page height for full capture
        total_height = self.driver.execute_script("return document.body.scrollHeight")
        self.driver.set_window_size(self.width, total_height)
        time.sleep(1)
        
        # Capture screenshot
        screenshot = self.driver.get_screenshot_as_png()
        return Image.open(io.BytesIO(screenshot))
    
    def split_into_slides(self, image, slide_height=1080):
        """Split long image into multiple slides"""
        slides = []
        img_width, img_height = image.size
        
        num_slides = (img_height + slide_height - 1) // slide_height
        
        for i in range(num_slides):
            top = i * slide_height
            bottom = min((i + 1) * slide_height, img_height)
            
            slide_img = image.crop((0, top, img_width, bottom))
            slides.append(slide_img)
            
        return slides
    
    def create_pptx(self, slides, output_path):
        """Create PPTX from slide images"""
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        
        # Remove default layouts
        blank_layout = prs.slide_layouts[6]
        
        for slide_img in slides:
            slide = prs.slides.add_slide(blank_layout)
            
            # Save image temporarily
            temp_img_path = 'temp_slide.png'
            slide_img.save(temp_img_path)
            
            # Add image to slide
            left = Inches(0)
            top = Inches(0)
            slide.shapes.add_picture(temp_img_path, left, top, 
                                    width=prs.slide_width, 
                                    height=prs.slide_height)
            
            # Clean up temp file
            os.remove(temp_img_path)
        
        prs.save(output_path)
        print(f"✅ PPTX created: {output_path}")
    
    def convert(self, html_path, output_path=None, slides_per_page=1):
        """Main conversion function"""
        if not output_path:
            output_path = html_path.replace('.html', '.pptx')
        
        print(f"🔄 Converting {html_path} to PPTX...")
        
        # Capture HTML
        print("📸 Capturing HTML rendering...")
        full_image = self.capture_html(html_path)
        
        # Split into slides
        print("✂️  Splitting into slides...")
        slides = self.split_into_slides(full_image, slide_height=self.height)
        print(f"📊 Created {len(slides)} slides")
        
        # Create PPTX
        print("📝 Generating PPTX...")
        self.create_pptx(slides, output_path)
        
        return output_path
    
    def cleanup(self):
        """Close browser"""
        if self.driver:
            self.driver.quit()

def main():
    if len(sys.argv) < 2:
        print("Usage: python html_to_pptx.py <input.html> [output.pptx]")
        sys.exit(1)
    
    html_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else None
    
    if not os.path.exists(html_file):
        print(f"❌ Error: File '{html_file}' not found")
        sys.exit(1)
    
    converter = HTMLtoPPTXConverter()
    try:
        converter.convert(html_file, output_file)
    finally:
        converter.cleanup()

if __name__ == "__main__":
    main()
